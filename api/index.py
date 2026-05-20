# =========================================================
# RESUME TO PPT BACKEND — HYBRID THEME
# Theme: Deloitte white/green/teal  |  Layout: Sidebar + Green Bars
# =========================================================

import fitz
import os
import re
import json
import tempfile
import traceback

from groq import Groq
from dotenv import load_dotenv

import pytesseract
from PIL import Image
from docx import Document

from flask import (
    Flask, request, send_file, jsonify, after_this_request
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# =========================================================
# CONFIG
# =========================================================

load_dotenv()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))
MODEL_NAME = "llama-3.3-70b-versatile"

# =========================================================
# DESIGN TOKENS  (WHITE / CHARTREUSE / TEAL HYBRID)
# =========================================================

PRIMARY    = RGBColor(140, 198, 62)     # Chartreuse green — header & section bars
ACCENT     = RGBColor(56,  118, 137)    # Teal — sidebar labels, accents
DARK       = RGBColor(30,  30,  30)     # Near-black — body headings
WHITE      = RGBColor(255, 255, 255)   # White
OFF_WHITE  = RGBColor(255, 255, 255)   # Slide background
SIDEBAR_BG = RGBColor(230, 242, 245)   # Very light teal-blue — left sidebar
PANEL_BG   = RGBColor(248, 248, 248)   # Light gray — right content backdrop
BORDER     = RGBColor(200, 220, 225)   # Light teal-gray — borders
TEXT       = RGBColor(50,  50,  50)    # Dark charcoal — body text
MUTED      = RGBColor(100, 110, 110)   # Medium gray — meta text
TAG_BG     = RGBColor(210, 235, 190)   # Light green — tag chips
TAG_TEXT   = RGBColor(30,  70,  20)    # Dark forest green — tag text

LIGHT_GREEN = RGBColor(230, 245, 210)   # Very light chartreuse — sidebar bottom
FONT_NAME = "Calibri"

# Slide dimensions
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Layout
MARGIN    = Inches(0.22)
HEADER_H  = Inches(0.85)
LEFT_W    = Inches(3.5)
GUTTER    = Inches(0.15)
RIGHT_X   = LEFT_W + MARGIN + GUTTER
RIGHT_W   = SLIDE_W - RIGHT_X - MARGIN

# =========================================================
# APP
# =========================================================

app = Flask(__name__)

# =========================================================
# TEXT HELPERS
# =========================================================

def clean_text(text):
    if not text:
        return ""
    text = re.sub(r'[\*_`#]', '', text)
    text = text.replace("\x00", " ")
    lines = []
    for line in text.splitlines():
        line = re.sub(r'\s+', ' ', line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def clean_filename(name):
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    return name.strip().replace(" ", "_")[:50] or "Resume"


def safe_str(val, default=""):
    return str(val).strip() if val else default


def clamp_list(lst, n):
    return lst[:n] if isinstance(lst, list) else []

# =========================================================
# PDF / DOCX EXTRACTION
# =========================================================

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    full_text = []
    for page in doc:
        blocks = page.get_text("blocks")
        blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
        page_text = "\n".join(
            block[4] for block in blocks if block[4].strip()
        )
        full_text.append(page_text)
    return "\n".join(full_text)


def ocr_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.close()
        pix.save(tmp.name)
        text += pytesseract.image_to_string(Image.open(tmp.name))
        try:
            os.remove(tmp.name)
        except:
            pass
    return text


def extract_text_from_docx(path):
    from docx.oxml.ns import qn

    doc = Document(path)
    full_text = []

    processed_table_ids = set()

    def is_education_like_table(table):
        if not table.rows or len(table.columns) < 2:
            return False
        first_cells = [c.text.strip() for c in table.rows[0].cells]
        if not any(first_cells):
            return False
        header_score = sum(
            1 for c in first_cells
            if c and len(c.split()) <= 6 and (c.istitle() or c.isupper() or c[0].isupper())
        )
        return header_score >= 2

    def table_to_markdown(table):
        rows_text = []
        for i, row in enumerate(table.rows):
            cells = []
            prev = None
            for cell in row.cells:
                txt = " ".join(cell.text.split())
                if txt != prev:
                    cells.append(txt)
                    prev = txt
            rows_text.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_text.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(rows_text)

    def table_to_flat(table):
        lines = []
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt:
                    lines.append(txt)
        return "\n".join(lines)

    def iter_block_items(parent):
        body = parent.element.body
        for child in body:
            if child.tag == qn('w:p'):
                yield ('para', child)
            elif child.tag == qn('w:tbl'):
                yield ('table', child)

    para_map  = {p._element: p for p in doc.paragraphs}
    table_map = {t._element: t for t in doc.tables}

    for kind, elem in iter_block_items(doc):
        if kind == 'para':
            para = para_map.get(elem)
            if para:
                t = para.text.strip()
                if t:
                    full_text.append(t)

        elif kind == 'table':
            table = table_map.get(elem)
            if table is None or id(elem) in processed_table_ids:
                continue
            processed_table_ids.add(id(elem))

            for nested_tbl_elem in elem.iter(qn('w:tbl')):
                if nested_tbl_elem is not elem:
                    processed_table_ids.add(id(nested_tbl_elem))

            if is_education_like_table(table):
                full_text.append(table_to_markdown(table))
            else:
                flat = table_to_flat(table)
                if flat:
                    full_text.append(flat)

    return "\n".join(full_text)


def extract_images_from_pdf(path):
    images = []
    try:
        doc = fitz.open(path)
        for page in doc:
            for img in page.get_images():
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha < 4:
                    p = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
                    pix.save(p)
                    images.append(p)
    except Exception as e:
        print("Image extraction error:", e)
    return images


def extract_images_from_docx(path):
    images = []
    try:
        doc = Document(path)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                p = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
                with open(p, "wb") as f:
                    f.write(img_data)
                images.append(p)
    except Exception as e:
        print("DOCX image error:", e)
    return images

# =========================================================
# AI — JSON EXTRACTION & VALIDATION
# =========================================================

def extract_json(content):
    content = re.sub(r"```json|```", "", content).strip()
    start = content.find("{")
    end = content.rfind("}")
    if start == -1 or end == -1:
        raise Exception("AI did not return valid JSON")
    return json.loads(content[start:end + 1])


def validate_ai_data(data):
    validated_exp = []
    for exp in clamp_list(data.get("experience", []), 4):
        validated_exp.append({
            "title":   safe_str(exp.get("title"),   "Role"),
            "company": safe_str(exp.get("company"), "Company"),
            "dates":   safe_str(exp.get("dates"),   ""),
            "summary": safe_str(exp.get("summary"), ""),
        })

    validated_projects = []
    for proj in clamp_list(data.get("projects", []), 3):
        validated_projects.append({
            "name":        safe_str(proj.get("name"),        "Project"),
            "tech":        safe_str(proj.get("tech"),        ""),
            "description": safe_str(proj.get("description"), ""),
        })

    raw_strengths = clamp_list(data.get("core_strengths", []), 6)
    strengths = []
    for s in raw_strengths:
        s = safe_str(s)
        words = s.split()
        tag = " ".join(words[:3]) if len(words) > 3 else s
        if tag:
            strengths.append(tag)

    certs = [safe_str(c) for c in clamp_list(data.get("certifications", []), 3) if safe_str(c)]

    raw_edu = data.get("education", "")
    if isinstance(raw_edu, list):
        edu_entries = [safe_str(e) for e in raw_edu if safe_str(e)]
    else:
        edu_entries = [safe_str(raw_edu)] if safe_str(raw_edu) else []

    raw_skills = clamp_list(data.get("skills", []), 20)
    skills = []
    for sk in raw_skills:
        sk = safe_str(sk)
        words = sk.split()
        tag = " ".join(words[:3]) if len(words) > 3 else sk
        if tag:
            skills.append(tag)

    return {
        "name":           safe_str(data.get("name"),        "Candidate"),
        "designation":    safe_str(data.get("designation"), "Professional"),
        "location":       safe_str(data.get("location"),    ""),
        "summary":        safe_str(data.get("summary"),     ""),
        "contact":        safe_str(data.get("contact"),     ""),
        "skills":         skills,
        "core_strengths": strengths,
        "education":      edu_entries,
        "certifications": certs,
        "projects":       validated_projects,
        "experience":     validated_exp,
    }

# =========================================================
# AI PROMPT
# =========================================================

def generate_summary_json(text):
    prompt = f"""
You are an elite resume parser. Extract structured data for a professional PowerPoint slide.

====================================================
STRICT OUTPUT RULES
====================================================
1. Return STRICT JSON only. Zero markdown, zero explanation, zero preamble.
2. NEVER truncate any field — return complete text.
3. NEVER copy experience bullets into projects. Projects must be real project names.
4. core_strengths must be SHORT TAGS: 2–3 words max (e.g. "API Design", "Team Leadership", "CI/CD"). NOT sentences.
5. skills must be a JSON array of individual skill strings — each skill name 1–3 words only (e.g. "Spring Boot"). NOT a comma-separated string.
6. experience.dates: extract real date range (e.g. "Jan 2021 – Mar 2023"). Use "" if not found.
7. certifications: short strings, one per item.
8. Do NOT hallucinate any information not in the resume.

====================================================
CRITICAL FIELD INSTRUCTIONS
====================================================

FIELD: summary
  - Write a SYNTHESISED, crisp 2-sentences professional bio of the candidate.
  - Read the ENTIRE resume (all experience, skills, achievements) and distil
    the candidate's identity into two punchy sentences.
  - Do NOT copy the first line of the resume's summary section.
  - Do NOT start with "Experienced" or "Seasoned" — find a more specific opener.
  - Target length: approximately 25 words (acceptable range: 20–30 words).
  - Example: "Frontend architect with 8+ years building scalable Angular
    enterprise apps, specialising in WCAG compliance and large-scale migrations."

FIELD: experience[].summary
  - For each role, write ONE crisp sentence that captures the KEY contribution
    of that role based on ALL bullet points and responsibilities listed in the resume.
  - Synthesise — do NOT copy a bullet verbatim. Do NOT use bullet-point style.
  - Must fit comfortably in ONE line (~15–20 words max).
  - Example: "Led Angular 11→15 migration and WCAG compliance for enterprise
    risk modules, cutting production issues through rigorous code reviews."

FIELD: projects[].description
  - One crisp sentence (~15 words max) describing what the project does / its impact.
  - Synthesise from all detail given in the resume for that project.
  - Do NOT copy resume text verbatim.

FIELD: projects[].tech
  - Comma-separated tech stack string, e.g. "React, Flask, Groq API".
  - 1–3 words per technology name.

FIELD: education
  - The resume text may contain a MARKDOWN TABLE for education, e.g.:
      | Qualification | Institution | University | Year of passing | Percentage |
      | --- | --- | --- | --- | --- |
      | MCA | REVA University | REVA University | 2020 | 78.5% |
      | BCA | REVA University | REVA University | 2018 | 75.6% |
  - If a table is present, read EVERY row and produce an "education" array.
  - If education is given as plain text, produce a single-item array.
  - Each entry: "Degree/Qualification, Institution (Year) — Percentage/Grade if present"
  - Examples of well-formatted entries:
      "MCA, REVA University (2020) — 78.5%"
      "B.Tech Computer Science, IIT Delhi (2019)"
      "12th, KV AFS Yelahanka / CBSE (2015) — 65%"
  - Keep entries concise — one line each, no extra narrative.

====================================================
REQUIRED JSON SCHEMA
====================================================
{{
  "name": "Full Name",
  "designation": "Current Job Title",
  "location": "City, Country",
  "summary": "<synthesised 1-sentence bio — see instructions above>",
  "contact": "email | phone",
  "skills": ["Skill1", "Skill2", "Skill3"],
  "core_strengths": ["Tag One", "Tag Two", "Tag Three"],
  "education": ["Degree/Qual, Institution (Year) — Grade", "Next Degree, Institution (Year)"],
  "certifications": ["Cert 1", "Cert 2"],
  "experience": [
    {{
      "title": "Role Title",
      "company": "Company Name",
      "dates": "Mon YYYY – Mon YYYY",
      "summary": "<one crisp sentence summarising ALL contributions in this role>"
    }}
  ],
  "projects": [
    {{
      "name": "Project Name",
      "tech": "Tech1, Tech2, Tech3",
      "description": "<one crisp sentence about what the project does or achieved>"
    }}
  ]
}}

====================================================
RESUME TEXT
====================================================
{text[:9000]}
"""

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.05,
        max_tokens=2500,
    )

    raw = response.choices[0].message.content
    print("\n===== AI RESPONSE =====\n", raw, "\n=======================\n")
    data = extract_json(raw)
    return validate_ai_data(data)

# =========================================================
# PPT DRAW PRIMITIVES
# =========================================================

def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    """Add a rectangle. radius=True uses rounded rectangle (shape type 5).
       radius=False uses standard rectangle (shape type 1) for sharper corners."""
    shape_type = 5 if radius else 1
    sp = slide.shapes.add_shape(shape_type, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(0.5)
    else:
        sp.line.fill.background()
    return sp


def add_textbox(
    slide, text, left, top, width, height,
    font_size=10, bold=False, color=TEXT,
    align=PP_ALIGN.LEFT, italic=False,
    word_wrap=True, font_name=None
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align

    run = p.runs[0] if p.runs else p.add_run()
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = font_name or FONT_NAME
    run.font.color.rgb = color

    return box


def sidebar_section_label(slide, title, left, top, width):
    """Teal uppercase label with a thin underline for the left sidebar."""
    add_textbox(
        slide, title.upper(),
        left, top, width, Inches(0.26),
        font_size=12, bold=True, color=ACCENT
    )
    line_y = top + Inches(0.24)
    add_rect(slide, left, line_y, width * 0.30, Inches(0.015), ACCENT)
    return line_y + Inches(0.08)


def section_header_bar(slide, title, left, top, width):
    """Chartreuse green horizontal bar with white uppercase text."""
    bar_h = Inches(0.32)
    add_rect(slide, left, top, width, bar_h, PRIMARY)
    add_textbox(
        slide, title.upper(),
        left + Inches(0.08), top + Inches(0.03),
        width - Inches(0.16), bar_h - Inches(0.06),
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT
    )
    return top + bar_h


def estimate_text_height(text, font_size_pt, box_width_inches, line_spacing=1.3):
    if not text:
        return 0
    avg_char_width_in = font_size_pt * 0.009
    chars_per_line = max(1, int(box_width_inches / avg_char_width_in))
    line_height_in = (font_size_pt / 72.0) * line_spacing
    words = str(text).split()
    lines = 1
    current = 0
    for w in words:
        wlen = len(w) + 1
        if current + wlen > chars_per_line:
            lines += 1
            current = wlen
        else:
            current += wlen
    return lines * line_height_in

# =========================================================
# SECTION RENDERERS
# =========================================================

def render_skills_tags_compact(slide, skills, left, top, panel_width_emu):
    """Compact skills renderer - smaller tags, tighter spacing."""
    if not skills:
        return top

    panel_width_in = panel_width_emu / 914400

    TAG_H     = Inches(0.22)
    TAG_PAD_Y = Inches(0.03)
    TAG_FONT  = 10
    row_h     = TAG_H + TAG_PAD_Y

    x = left + Inches(0.03)
    y = top
    right_limit = left + panel_width_emu - Inches(0.06)

    for skill in skills:
        tag_w = Inches(len(skill) * 0.094 + 0.22)
        tag_w = min(tag_w, panel_width_emu - Inches(0.15))

        if x + tag_w > right_limit:
            x = left + Inches(0.03)
            y += row_h

        add_rect(slide, x, y, tag_w, TAG_H, TAG_BG, line_color=ACCENT)
        add_textbox(
            slide, skill,
            x + Inches(0.06), y + Inches(0.02),
            tag_w - Inches(0.08), TAG_H - Inches(0.03),
            font_size=TAG_FONT, color=TAG_TEXT, bold=True
        )
        x += tag_w + Inches(0.04)

    return y + row_h + Inches(0.02)


def render_experience(slide, experience, left, top, width, bottom_limit):
    y = top
    for i, exp in enumerate(experience):
        if i > 0:
            y += Inches(0.08)

        if y + Inches(0.50) > bottom_limit:
            break

        role_text    = exp["title"]
        company_text = exp["company"]
        dates_text   = exp.get("dates", "")
        summary_text = exp.get("summary", "")

        add_textbox(
            slide, role_text,
            left, y, width * 0.7, Inches(0.28),
            font_size=13, bold=True, color=DARK
        )

        meta = f"{company_text}  •  {dates_text}" if dates_text else company_text
        add_textbox(
            slide, meta,
            left, y + Inches(0.26), width, Inches(0.21),
            font_size=11, bold=True, color=MUTED, italic=True
        )
        y += Inches(0.50)

        if summary_text and y + Inches(0.20) <= bottom_limit:
            width_in = width / 914400
            est_h = estimate_text_height(summary_text, 11, width_in - 0.15) + 0.05
            est_h_emu = Inches(max(est_h, 0.22))
            if y + est_h_emu > bottom_limit:
                est_h_emu = bottom_limit - y - Inches(0.02)
            if est_h_emu > 0:
                add_textbox(
                    slide, f"▸  {summary_text}",
                    left + Inches(0.08), y,
                    width - Inches(0.12), est_h_emu,
                    font_size=11, bold=True, color=TEXT
                )
                y += est_h_emu + Inches(0.04)

        y += Inches(0.03)

    return y


def render_projects(slide, projects, left, top, width, bottom_limit):
    y = top
    for proj in projects:
        if y + Inches(0.35) > bottom_limit:
            break

        add_textbox(
            slide, proj["name"],
            left, y, width, Inches(0.28),
            font_size=13, bold=True, color=DARK
        )
        y += Inches(0.28)

        if proj.get("tech") and y + Inches(0.20) < bottom_limit:
            add_textbox(
                slide, f"Tech: {proj['tech']}",
                left, y, width, Inches(0.20),
                font_size=10, bold=True, color=MUTED, italic=True
            )
            y += Inches(0.21)

        description = proj.get("description", "")
        if description and y + Inches(0.20) <= bottom_limit:
            width_in = width / 914400
            est_h = estimate_text_height(description, 11, width_in - 0.15) + 0.05
            est_h_emu = Inches(max(est_h, 0.22))
            if y + est_h_emu > bottom_limit:
                est_h_emu = bottom_limit - y - Inches(0.02)
            if est_h_emu > 0:
                add_textbox(
                    slide, f"▸  {description}",
                    left + Inches(0.08), y,
                    width - Inches(0.12), est_h_emu,
                    font_size=11, bold=True, color=TEXT
                )
                y += est_h_emu + Inches(0.04)

        y += Inches(0.05)

    return y

# =========================================================
# MAIN PPT BUILDER — HYBRID TABLE/TEXT LAYOUT
# =========================================================

DEFAULT_AVATAR_PATH = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "demo_avatar.png")
)

def create_resume_ppt(data, images=None):

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ---- 1. SLIDE BACKGROUND ----
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)

    has_resume_photo = bool(images and len(images) > 0)

    if not has_resume_photo and os.path.exists(DEFAULT_AVATAR_PATH):
        images = [DEFAULT_AVATAR_PATH]
        has_photo = True
    else:
        has_photo = has_resume_photo

    # =========================================================
    # LAYOUT: FULL-WIDTH TOP HEADER + LEFT SIDEBAR + RIGHT PANEL
    # =========================================================

    # ---- 2. TOP HEADER BAR (FULL WIDTH with PADDING) ----
    HEADER_PAD = Inches(0.08)
    header_h = Inches(0.90)
    header_top = HEADER_PAD
    header_left = HEADER_PAD
    header_width = SLIDE_W - (HEADER_PAD * 2)

    add_rect(slide, header_left, header_top, header_width, header_h, PRIMARY)

    content_pad = Inches(0.15)
    left_content_x = header_left + content_pad

    if has_photo and images and len(images) > 0:
        photo_size = Inches(0.65)
        photo_left = left_content_x
        photo_top = header_top + (header_h - photo_size) / 2
        try:
            slide.shapes.add_picture(
                images[0], photo_left, photo_top,
                width=photo_size, height=photo_size
            )
            name_left = photo_left + photo_size + Inches(0.12)
        except Exception as e:
            print("Photo insert error:", e)
            name_left = left_content_x
    else:
        name_left = left_content_x

    name_width = Inches(4.5)
    name_top = header_top + Inches(0.10)

    add_textbox(
        slide, data["name"].upper(),
        name_left, name_top,
        name_width, Inches(0.40),
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT
    )

    add_textbox(
        slide, data["designation"],
        name_left, name_top + Inches(0.40),
        name_width, Inches(0.26),
        font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT, italic=True
    )

    contact_parts = []
    if data["contact"]:
        contact_parts.extend([p.strip() for p in data["contact"].split("|")])
    if data["location"]:
        contact_parts.append(data["location"].strip())

    if contact_parts:
        contact_line = " || ".join(contact_parts)
        contact_width = Inches(5.0)
        contact_left = header_left + header_width - contact_width - content_pad 
        contact_top = header_top + (header_h - Inches(0.30)) / 2

        add_textbox(
            slide, contact_line,
            contact_left, contact_top,
            contact_width, Inches(0.30),
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT
        )

    # ---- 3. LEFT SIDEBAR (below header, with INCREASED top spacing) ----
    sidebar_top = header_top + header_h + Inches(0.18)
    sidebar_h = SLIDE_H - sidebar_top - HEADER_PAD
    sidebar_left = HEADER_PAD
    sidebar_width = LEFT_W

    add_rect(slide, sidebar_left, sidebar_top, sidebar_width, sidebar_h, LIGHT_GREEN, line_color=BORDER)

    sx = sidebar_left + Inches(0.18)
    sw = sidebar_width - Inches(0.36)
    sy = sidebar_top + Inches(0.14)
    sidebar_bottom = sidebar_top + sidebar_h - Inches(0.10)

    # -- SUMMARY (in left sidebar) --
    if data["summary"] and sy < sidebar_bottom - Inches(0.5):
        sy = sidebar_section_label(slide, "Summary", sx, sy, sw)
        summary_w_in = sw / 914400
        summary_h = estimate_text_height(data["summary"], 11, summary_w_in) + 0.05
        summary_h = max(summary_h, 0.25)
        summary_h_emu = Inches(min(summary_h, 1.0))
        if sy + summary_h_emu <= sidebar_bottom:
            add_textbox(
                slide, data["summary"], sx, sy, sw, summary_h_emu,
                font_size=11, bold=True, color=TEXT , align=PP_ALIGN.JUSTIFY
            )
            sy += summary_h_emu + Inches(0.08)

    # -- KEY SKILLS (COMPACT) --
    if data["skills"] and sy < sidebar_bottom - Inches(0.5):
        sy = sidebar_section_label(slide, "Key Skills", sx, sy, sw)
        sy = render_skills_tags_compact(slide, data["skills"], sx, sy, int(sw))
        sy += Inches(0.08)

    # -- CORE STRENGTHS (COMPACT) --
    if data["core_strengths"] and sy < sidebar_bottom - Inches(0.5):
        sy = sidebar_section_label(slide, "Core Strengths", sx, sy, sw)
        sy = render_skills_tags_compact(slide, data["core_strengths"], sx, sy, int(sw))
        sy += Inches(0.08)

    # -- EDUCATION --
    if data["education"] and sy < sidebar_bottom - Inches(0.5):
        sy = sidebar_section_label(slide, "Education", sx, sy, sw)
        for edu_entry in data["education"]:
            if sy > sidebar_bottom - Inches(0.30):
                break
            entry_h = Inches(max(estimate_text_height(edu_entry, 11, sw / 914400) + 0.04, 0.25))
            add_textbox(slide, f"• {edu_entry}", sx, sy, sw, entry_h,
                       font_size=11, bold=True, color=TEXT)
            sy += entry_h + Inches(0.04)
        sy += Inches(0.08)

    # -- CERTIFICATIONS --
    if data["certifications"] and sy < sidebar_bottom - Inches(0.5):
        sy = sidebar_section_label(slide, "Certifications", sx, sy, sw)
        for cert in data["certifications"]:
            if sy > sidebar_bottom - Inches(0.30):
                break
            add_textbox(slide, f"• {cert}", sx, sy, sw, Inches(0.24),
                       font_size=11, bold=True, color=TEXT)
            sy += Inches(0.25)
        sy += Inches(0.08)

    # ---- 4. RIGHT PANEL (below header, with INCREASED top spacing) ----
    rp_left = sidebar_left + sidebar_width + GUTTER
    rp_w = SLIDE_W - rp_left - HEADER_PAD
    rp_top = sidebar_top
    rp_h = sidebar_h
    rp_bottom = rp_top + rp_h - Inches(0.08)

    add_rect(
        slide,
        rp_left, rp_top,
        rp_w, rp_h,
        PANEL_BG, line_color=BORDER
    )

    content_pad_x = Inches(0.08)
    content_pad_y = Inches(0.06)
    ry = rp_top + content_pad_y
    rx = rp_left + content_pad_x
    rw = rp_w - (content_pad_x * 2)

    # -- HYBRID LOGIC: Determine table vs text format --
    exp_count = len(data.get("experience", []))
    proj_count = len(data.get("projects", []))
    total_items = exp_count + proj_count

    use_exp_table = False
    use_proj_table = False

    if total_items > 4:
        if exp_count > 2 and proj_count > 2:
            use_exp_table = True
            use_proj_table = True
        elif exp_count > 2:
            use_exp_table = True
            use_proj_table = False
        elif proj_count > 2:
            use_exp_table = False
            use_proj_table = True

    # -- EXPERIENCE --
    if data["experience"]:
        if use_exp_table:
            ry = render_experience_table(slide, data["experience"], rx, ry, rw, rp_bottom)
        else:
            bar_bottom = section_header_bar(slide, "Experience", rx, ry, rw)
            ry = bar_bottom + Inches(0.10)

            if data["projects"] and not use_proj_table:
                proj_reserve = Inches(0.35) + Inches(0.10) + Inches(0.75)
                exp_bottom = rp_bottom - proj_reserve - Inches(0.15)
            else:
                exp_bottom = rp_bottom

            ry = render_experience(slide, data["experience"], rx, ry, rw, bottom_limit=exp_bottom)
            ry += Inches(0.10)

    # -- PROJECTS --
    if data["projects"]:
        if use_proj_table:
            ry = render_projects_table(slide, data["projects"], rx, ry, rw, rp_bottom)
        else:
            bar_bottom = section_header_bar(slide, "Projects", rx, ry, rw)
            ry = bar_bottom + Inches(0.10)
            render_projects(slide, data["projects"], rx, ry, rw, bottom_limit=rp_bottom)

    # =========================================================
    # SAVE
    # =========================================================

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.close()
    prs.save(tmp.name)
    return tmp.name


# =========================================================
# TABLE RENDERERS — Used when >2 items
# =========================================================

def render_experience_table(slide, experience, left, top, width, bottom_limit):
    """Render experience as a table when >2 entries."""
    bar_h = Inches(0.35)
    add_rect(slide, left, top, width, bar_h, PRIMARY)
    add_textbox(
        slide, "EXPERIENCE",
        left + Inches(0.08), top + Inches(0.04),
        width - Inches(0.16), bar_h - Inches(0.08),
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT
    )

    y = top + bar_h + Inches(0.08)

    col_widths = [width * 0.22, width * 0.20, width * 0.15, width * 0.43]
    row_h = Inches(0.42)

    headers = ["Role", "Company", "Duration", "Summary"]
    x_offset = left
    for i, header in enumerate(headers):
        add_rect(slide, x_offset, y, col_widths[i], row_h, PRIMARY)
        add_textbox(slide, header, x_offset + Inches(0.04), y + Inches(0.10),
                   col_widths[i] - Inches(0.08), row_h - Inches(0.20),
                   font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x_offset += col_widths[i]

    y += row_h

    for idx, exp in enumerate(experience):
        if y + row_h > bottom_limit - Inches(0.15):
            break

        role = exp.get("title", "")
        company = exp.get("company", "")
        dates = exp.get("dates", "")
        summary = exp.get("summary", "")

        row_bg = WHITE if idx % 2 == 0 else RGBColor(250, 250, 250)
        add_rect(slide, left, y, width, row_h, row_bg, line_color=BORDER)

        x_offset = left
        for i in range(4):
            x_offset += col_widths[i]
            add_rect(slide, x_offset - Inches(0.01), y, Inches(0.01), row_h, BORDER)
        add_rect(slide, left, y + row_h - Inches(0.01), width, Inches(0.01), BORDER)

        cells = [role, company, dates, summary]
        x_offset = left
        for i, text in enumerate(cells):
            add_textbox(slide, str(text), x_offset + Inches(0.04), y + Inches(0.06),
                       col_widths[i] - Inches(0.08), row_h - Inches(0.12),
                       font_size=11, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
            x_offset += col_widths[i]

        y += row_h

    return y + Inches(0.10)


def render_projects_table(slide, projects, left, top, width, bottom_limit):
    """Render projects as a table when >2 entries."""
    bar_h = Inches(0.35)
    add_rect(slide, left, top, width, bar_h, PRIMARY)
    add_textbox(
        slide, "PROJECTS",
        left + Inches(0.08), top + Inches(0.04),
        width - Inches(0.16), bar_h - Inches(0.08),
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT
    )

    y = top + bar_h + Inches(0.08)

    col_widths = [width * 0.30, width * 0.25, width * 0.45]
    row_h = Inches(0.42)

    headers = ["Project", "Tech / Tools", "Outcome"]
    x_offset = left
    for i, header in enumerate(headers):
        add_rect(slide, x_offset, y, col_widths[i], row_h, PRIMARY)
        add_textbox(slide, header, x_offset + Inches(0.04), y + Inches(0.10),
                   col_widths[i] - Inches(0.08), row_h - Inches(0.20),
                   font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x_offset += col_widths[i]

    y += row_h

    for idx, proj in enumerate(projects):
        if y + row_h > bottom_limit - Inches(0.15):
            break

        name = proj.get("name", "")
        tech = proj.get("tech", "")
        desc = proj.get("description", "")

        row_bg = WHITE if idx % 2 == 0 else RGBColor(250, 250, 250)
        add_rect(slide, left, y, width, row_h, row_bg, line_color=BORDER)

        x_offset = left
        for i in range(3):
            x_offset += col_widths[i]
            add_rect(slide, x_offset - Inches(0.01), y, Inches(0.01), row_h, BORDER)
        add_rect(slide, left, y + row_h - Inches(0.01), width, Inches(0.01), BORDER)

        cells = [name, tech, desc]
        x_offset = left
        for i, text in enumerate(cells):
            add_textbox(slide, str(text), x_offset + Inches(0.04), y + Inches(0.06),
                       col_widths[i] - Inches(0.08), row_h - Inches(0.12),
                       font_size=11, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
            x_offset += col_widths[i]

        y += row_h

    return y + Inches(0.10)
# =========================================================
# CORS
# =========================================================

@app.after_request
def add_cors_headers(response):
    response.headers["Access-Control-Allow-Origin"]  = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type,Authorization"
    response.headers["Access-Control-Allow-Methods"] = "GET,POST,OPTIONS"
    response.headers["Access-Control-Expose-Headers"] = "Content-Disposition"
    return response

# =========================================================
# ROUTE
# =========================================================

@app.route("/generate-ppt", methods=["POST", "OPTIONS"])
def generate_ppt_api():

    if request.method == "OPTIONS":
        return jsonify({"status": "ok"}), 200

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file   = request.files["file"]
    suffix = os.path.splitext(file.filename)[1].lower()

    tmp_in = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp_in.close()
    file.save(tmp_in.name)

    images = []

    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(tmp_in.name)
            if len(text.strip()) < 400:
                text = ocr_pdf(tmp_in.name)
            images = extract_images_from_pdf(tmp_in.name)

        elif suffix == ".docx":
            text   = extract_text_from_docx(tmp_in.name)
            images = extract_images_from_docx(tmp_in.name)

        else:
            return jsonify({"error": "Only PDF and DOCX are supported"}), 400

        text     = clean_text(text)
        data     = generate_summary_json(text)
        ppt_path = create_resume_ppt(data, images)
        filename = clean_filename(data["name"]) + "_Resume.pptx"

        @after_this_request
        def cleanup(response):
            for p in [tmp_in.name, ppt_path] + images:
                try:
                    os.remove(p)
                except:
                    pass
            return response

        return send_file(
            ppt_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        traceback.print_exc()
        try:
            os.remove(tmp_in.name)
        except:
            pass
        for p in images:
            try:
                os.remove(p)
            except:
                pass
        return jsonify({"error": str(e)}), 500

# =========================================================
# MAIN
# =========================================================

if __name__ == "__main__":
    app.run(debug=True, port=5000)